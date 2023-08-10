# This class will handle parsing the presentation file and extracting its data.

from pptx import Presentation

class PresentationParser:
    def __init__(self, file_path):
        self.file_path = file_path

    def parse_presentation(self):
        # Logic to parse the presentation file and extract data
        prs = Presentation(self.file_path)
        return prs

    def extract_text(self, slide):
        # Logic to extract text from a slide
        text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text.strip() + " "
        return text.strip()